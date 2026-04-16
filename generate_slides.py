"""
Generate HCP AI Smart Inbox presentation deck (.pptx)
HouseCall Pro brand: Navy #08151E, Blue #0055FF, Gold #FFB706
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Brand constants ──────────────────────────────────────────────
NAVY   = RGBColor(0x08, 0x15, 0x1E)
BLUE   = RGBColor(0x00, 0x55, 0xFF)
GOLD   = RGBColor(0xFF, 0xB7, 0x06)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY  = RGBColor(0x32, 0x37, 0x3C)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
MGRAY  = RGBColor(0x8A, 0x8F, 0x98)
BLUE_LIGHT = RGBColor(0x1A, 0x6B, 0xFF)
BLUE_BG    = RGBColor(0x0A, 0x2A, 0x5E)
GOLD_DIM   = RGBColor(0xFF, 0xC8, 0x40)

FONT_HEADING = "Lato"
FONT_BODY    = "Lato"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────
def add_bg(slide, color):
    """Set solid background fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.2):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Add a text box and return the text frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=16, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_after=6, space_before=0,
             font_name=FONT_BODY, line_spacing=1.15):
    """Append a paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner rounding
    shape.adjustments[0] = 0.06
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow_right(slide, left, top, width, height, color):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, color):
    """Add a chevron (notched right arrow)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── SLIDE 1: "I Started by Listening" ───────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
add_bg(slide1, NAVY)

# Title
add_textbox(slide1, Inches(0.8), Inches(0.4), Inches(7), Inches(0.8),
            "I Started by Listening", font_size=40, color=WHITE, bold=True,
            font_name=FONT_HEADING)

# Subtitle
add_textbox(slide1, Inches(0.8), Inches(1.15), Inches(7), Inches(0.5),
            "Rauen Flooring  \u2022  Fort Myers, FL  \u2022  3 people",
            font_size=20, color=GOLD, bold=False)

# Photo placeholder
ph = add_rounded_rect(slide1, Inches(9.0), Inches(0.5), Inches(3.8), Inches(2.8),
                       RGBColor(0x1A, 0x25, 0x33), border_color=MGRAY, border_width=Pt(1))
add_textbox(slide1, Inches(9.0), Inches(1.5), Inches(3.8), Inches(0.8),
            "[ Family Photo ]", font_size=16, color=MGRAY,
            alignment=PP_ALIGN.CENTER)

# Quotes with blue accent bars
quotes = [
    ("\u201cI\u2019ll finish a job covered in dust, check my phone, and there are "
     "six messages. The one about the invoice? Buried.\u201d",
     "\u2014 Sandro"),
    ("\u201cCustomers always ask \u2018why did you charge for leveling?\u2019 "
     "We\u2019ve explained it a hundred times, but when we\u2019re on a job "
     "we can\u2019t stop and type it out.\u201d",
     "\u2014 Fabiana"),
    ("\u201cI thought Sandro answered her. He thought I did. She waited "
     "three days and hired someone else.\u201d",
     "\u2014 Fabiana"),
]

y_start = Inches(3.6)
for i, (quote, attr) in enumerate(quotes):
    y = y_start + Inches(i * 1.15)
    # Blue accent bar
    add_rect(slide1, Inches(0.8), y, Inches(0.06), Inches(0.85), BLUE)
    # Quote text
    tf = add_rich_textbox(slide1, Inches(1.1), y, Inches(7.5), Inches(0.85))
    add_para(tf, quote, font_size=14, color=WHITE, line_spacing=1.25)
    add_para(tf, attr, font_size=13, color=GOLD, bold=True, space_before=2)

# Bottom callout
add_textbox(slide1, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
            "They were also clear about one thing: as customers themselves, "
            "they don\u2019t want to talk to a bot. They want to know a person "
            "is on the other end.",
            font_size=13, color=MGRAY)


# ── SLIDE 2: "The Data Backs It Up" ─────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_bg(slide2, WHITE)

add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "The Data Backs It Up", font_size=40, color=NAVY, bold=True,
            font_name=FONT_HEADING)

add_textbox(slide2, Inches(0.8), Inches(1.15), Inches(10), Inches(0.5),
            "What I heard at the kitchen table is not just my family\u2019s problem.",
            font_size=18, color=DGRAY)

# Stat cards — 2x2 grid
stats = [
    ("56%", "of small businesses are owed\nmoney from unpaid invoices",
     "$17.5K", "average amount owed\nper business"),
    ("94% \u2192 50%", "collection probability drops\nfrom 30 days to 6 months",
     "+50%", "payment likelihood increase\nwith just one reminder"),
]

card_w = Inches(5.2)
card_h = Inches(1.9)
x_positions = [Inches(0.8), Inches(6.6)]
y_positions = [Inches(2.0), Inches(4.2)]

for row_i, row in enumerate(stats):
    for col_i in range(2):
        idx = col_i * 2
        number = row[idx]
        desc = row[idx + 1]
        x = x_positions[col_i]
        y = y_positions[row_i]

        # Card background
        add_rounded_rect(slide2, x, y, card_w, card_h, LGRAY)

        # Number
        add_textbox(slide2, x + Inches(0.4), y + Inches(0.25), card_w - Inches(0.8),
                    Inches(0.7), number, font_size=36, color=GOLD, bold=True,
                    font_name=FONT_HEADING)

        # Description
        tf = add_rich_textbox(slide2, x + Inches(0.4), y + Inches(1.0),
                              card_w - Inches(0.8), Inches(0.8))
        for line in desc.split("\n"):
            add_para(tf, line, font_size=14, color=DGRAY, line_spacing=1.3)

# Bottom insight
add_textbox(slide2, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
            "Speed and clarity are the two biggest levers for getting paid.",
            font_size=20, color=NAVY, bold=True, font_name=FONT_HEADING)


# ── SLIDE 3: "Where the Funnel Leaks" ───────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_bg(slide3, WHITE)

add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "Where the Funnel Leaks", font_size=40, color=NAVY, bold=True,
            font_name=FONT_HEADING)

# Funnel stages
funnel_labels = [
    "Invoice\nSent",
    "Customer\nReceives",
    "Questions &\nConcerns",
    "Pro\nResponds",
    "Payment\nMade",
]
funnel_colors = [BLUE, BLUE, GOLD, BLUE, BLUE]

stage_w = Inches(2.1)
stage_h = Inches(0.95)
gap = Inches(0.25)
funnel_x_start = Inches(0.8)
funnel_y = Inches(1.5)

for i, (label, color) in enumerate(zip(funnel_labels, funnel_colors)):
    x = funnel_x_start + (stage_w + gap) * i
    shape = add_chevron(slide3, x, funnel_y, stage_w, stage_h, color)
    # Text on chevron
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE if color == BLUE else NAVY
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

# Gap indicator — small arrow pointing to the gold stage
add_textbox(slide3, Inches(4.7), Inches(2.55), Inches(2.5), Inches(0.45),
            "\u2191 THE GAP", font_size=14, color=GOLD, bold=True,
            alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

# Divider line
add_rect(slide3, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.02), LGRAY)

# What HCP does well
tf_hcp = add_rich_textbox(slide3, Inches(0.8), Inches(3.4), Inches(5.4), Inches(1.6))
add_para(tf_hcp, "What HCP does well today", font_size=16, color=NAVY,
         bold=True, space_after=8)
add_para(tf_hcp, "\u2022  Invoice creation & multiple payment methods", font_size=13,
         color=DGRAY, space_after=3)
add_para(tf_hcp, "\u2022  Automated email reminders", font_size=13,
         color=DGRAY, space_after=3)
add_para(tf_hcp, "\u2022  Marketing AI for invoice message templates", font_size=13,
         color=DGRAY, space_after=3)

# Competitors
tf_comp = add_rich_textbox(slide3, Inches(6.8), Inches(3.4), Inches(5.8), Inches(1.6))
add_para(tf_comp, "What competitors are adding", font_size=16, color=NAVY,
         bold=True, space_after=8)
add_para(tf_comp, "ServiceTitan \u2014 AI invoice summaries + one-click collection emails",
         font_size=13, color=DGRAY, space_after=3)
add_para(tf_comp, "Jobber \u2014 AI rewrite for invoice messages, reminders, two-way texting",
         font_size=13, color=DGRAY, space_after=3)
add_para(tf_comp, "Workiz \u2014 AI messaging across channels, tailored to customer history",
         font_size=13, color=DGRAY, space_after=3)

# Opportunity statement
add_rounded_rect(slide3, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5),
                 RGBColor(0x08, 0x15, 0x1E))

tf_opp = add_rich_textbox(slide3, Inches(1.2), Inches(5.7), Inches(11.0), Inches(1.2))
add_para(tf_opp, "THE OPPORTUNITY", font_size=12, color=GOLD, bold=True,
         space_after=6)
add_para(tf_opp, "Competitors are adding AI to individual touchpoints. None have a unified "
         "AI layer that watches the full conversation, understands context from memory, "
         "and helps the pro respond at the moment it matters most.",
         font_size=15, color=WHITE, line_spacing=1.35)


# ── SLIDE 4: "AI Smart Inbox" ───────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_bg(slide4, NAVY)

add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "AI Smart Inbox", font_size=44, color=WHITE, bold=True,
            font_name=FONT_HEADING)

add_textbox(slide4, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
            "An AI co-worker that lives inside HouseCall Pro\u2019s messaging inbox.",
            font_size=20, color=GOLD)

# Three pillar cards
pillars = [
    ("PRIORITIZE", "Surface what needs your attention right now. "
     "No hunting through threads."),
    ("BE QUICK", "AI-drafted responses ready the moment you check your inbox. "
     "Turn days into minutes."),
    ("DRIVE CLARITY", "Use context that already exists \u2014 conversation history, "
     "invoice details, your own FAQ answers \u2014 to draft accurate, helpful replies."),
]

pillar_w = Inches(3.6)
pillar_h = Inches(2.8)
pillar_gap = Inches(0.4)
pillar_x_start = Inches(0.8)
pillar_y = Inches(2.3)

for i, (label, desc) in enumerate(pillars):
    x = pillar_x_start + (pillar_w + pillar_gap) * i

    # Card background
    card = add_rounded_rect(slide4, x, pillar_y, pillar_w, pillar_h,
                            RGBColor(0x0F, 0x1F, 0x2E))

    # Blue top border accent
    add_rect(slide4, x + Inches(0.3), pillar_y + Inches(0.05),
             pillar_w - Inches(0.6), Inches(0.05), BLUE)

    # Pillar label
    add_textbox(slide4, x + Inches(0.3), pillar_y + Inches(0.35),
                pillar_w - Inches(0.6), Inches(0.4),
                label, font_size=16, color=GOLD, bold=True,
                font_name=FONT_HEADING)

    # Pillar description
    add_textbox(slide4, x + Inches(0.3), pillar_y + Inches(0.9),
                pillar_w - Inches(0.6), Inches(1.8),
                desc, font_size=14, color=WHITE, line_spacing=1.4)

# Design principle
add_textbox(slide4, Inches(0.8), Inches(5.7), Inches(11), Inches(0.8),
            "Design principle: AI works behind the scenes. The pro reviews, edits, "
            "and sends every message. No customer ever talks to a bot.",
            font_size=15, color=MGRAY, line_spacing=1.3)


# ── SLIDE 5: "Under the Hood" ───────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_bg(slide5, WHITE)

add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "Under the Hood", font_size=40, color=NAVY, bold=True,
            font_name=FONT_HEADING)

add_textbox(slide5, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "Two-stage pipeline with a persistent memory layer",
            font_size=18, color=DGRAY)

# Evaluate box
eval_w = Inches(4.5)
eval_h = Inches(2.2)
eval_x = Inches(0.8)
eval_y = Inches(2.0)
add_rounded_rect(slide5, eval_x, eval_y, eval_w, eval_h, BLUE)

tf_eval = add_rich_textbox(slide5, eval_x + Inches(0.3), eval_y + Inches(0.2),
                           eval_w - Inches(0.6), eval_h - Inches(0.4))
add_para(tf_eval, "1  EVALUATE", font_size=18, color=WHITE, bold=True, space_after=10)
add_para(tf_eval, "Every incoming message runs through an AI agent that decides:",
         font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF), space_after=6)
add_para(tf_eval, "\u2022  Does this need action?", font_size=13, color=WHITE, space_after=3)
add_para(tf_eval, "\u2022  What kind? How urgent?", font_size=13, color=WHITE, space_after=3)
add_para(tf_eval, "\u2022  Update customer memory", font_size=13, color=WHITE, space_after=3)

# Arrow between boxes
add_arrow_right(slide5, Inches(5.55), Inches(2.7), Inches(0.7), Inches(0.5), BLUE)

# Generate box
gen_x = Inches(6.5)
gen_y = Inches(2.0)
gen_w = Inches(4.5)
gen_h = Inches(2.2)
add_rounded_rect(slide5, gen_x, gen_y, gen_w, gen_h, BLUE)

tf_gen = add_rich_textbox(slide5, gen_x + Inches(0.3), gen_y + Inches(0.2),
                          gen_w - Inches(0.6), gen_h - Inches(0.4))
add_para(tf_gen, "2  GENERATE", font_size=18, color=WHITE, bold=True, space_after=10)
add_para(tf_gen, "If action needed, a second pass drafts a response using:",
         font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF), space_after=6)
add_para(tf_gen, "\u2022  Invoice details", font_size=13, color=WHITE, space_after=3)
add_para(tf_gen, "\u2022  Conversation history", font_size=13, color=WHITE, space_after=3)
add_para(tf_gen, "\u2022  Two layers of memory", font_size=13, color=WHITE, space_after=3)

# Memory bar
mem_y = Inches(4.55)
mem_h = Inches(1.3)
add_rounded_rect(slide5, Inches(0.8), mem_y, Inches(10.2), mem_h, GOLD)

# Memory bar title
add_textbox(slide5, Inches(1.1), mem_y + Inches(0.1), Inches(9.6), Inches(0.4),
            "MEMORY LAYER", font_size=14, color=NAVY, bold=True,
            font_name=FONT_HEADING)

# Business memory
add_textbox(slide5, Inches(1.1), mem_y + Inches(0.5), Inches(4.5), Inches(0.7),
            "Business Knowledge\nYour services, pricing, FAQ answers. "
            "Things you explain every week.",
            font_size=12, color=RGBColor(0x3A, 0x2A, 0x00), line_spacing=1.3)

# Divider
add_rect(slide5, Inches(5.9), mem_y + Inches(0.35), Inches(0.02), Inches(0.7),
         RGBColor(0xD9, 0x9A, 0x00))

# Customer memory
add_textbox(slide5, Inches(6.2), mem_y + Inches(0.5), Inches(4.5), Inches(0.7),
            "Customer History\nPreferences, past jobs, communication style. "
            "Grows with every interaction.",
            font_size=12, color=RGBColor(0x3A, 0x2A, 0x00), line_spacing=1.3)

# Skills list
tf_skills = add_rich_textbox(slide5, Inches(0.8), Inches(6.15), Inches(11), Inches(0.5))
add_para(tf_skills, "Skills:  Follow-up  \u2022  Question Response  \u2022  "
         "Invoice Summary  \u2022  Payment Plan    "
         "\u2014  Adding a new skill means writing a new prompt, not rebuilding the system.",
         font_size=13, color=DGRAY)

# Architecture callout
add_textbox(slide5, Inches(0.8), Inches(6.7), Inches(11), Inches(0.5),
            "The same pipeline works for any home service vertical and "
            "can extend beyond invoicing into other HCP workflows.",
            font_size=14, color=NAVY, bold=True)


# ── SLIDE 6: "Trust is Earned" ──────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
add_bg(slide6, WHITE)

add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "Trust is Earned", font_size=40, color=NAVY, bold=True,
            font_name=FONT_HEADING)

add_textbox(slide6, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "Progression is earned by metrics, not by a calendar.",
            font_size=18, color=DGRAY)

# Phase cards with progression
phases = [
    ("PHASE 1", "This MVP", "Every draft reviewed\nby the pro",
     "Building confidence", BLUE),
    ("PHASE 2", "Next", "Auto-send for proven,\nlow-risk actions",
     "Track record earns automation", BLUE_LIGHT),
    ("PHASE 3", "Future", "Autonomous for\nroutine cases",
     "AI operates where trust\nis established", GOLD),
]

phase_w = Inches(3.4)
phase_h = Inches(2.5)
phase_gap = Inches(0.55)
phase_x_start = Inches(0.8)
phase_y = Inches(1.9)

for i, (phase_label, timing, desc, subtext, color) in enumerate(phases):
    x = phase_x_start + (phase_w + phase_gap) * i

    add_rounded_rect(slide6, x, phase_y, phase_w, phase_h, color)

    text_color = WHITE if color != GOLD else NAVY

    # Phase label
    add_textbox(slide6, x + Inches(0.3), phase_y + Inches(0.2),
                phase_w - Inches(0.6), Inches(0.35),
                f"{phase_label}  \u2022  {timing}",
                font_size=12, color=text_color, bold=True,
                font_name=FONT_HEADING)

    # Description
    add_textbox(slide6, x + Inches(0.3), phase_y + Inches(0.7),
                phase_w - Inches(0.6), Inches(0.9),
                desc, font_size=18, color=text_color, bold=True,
                line_spacing=1.25)

    # Subtext
    sub_color = RGBColor(0xCC, 0xDD, 0xFF) if color != GOLD else RGBColor(0x5A, 0x4A, 0x00)
    add_textbox(slide6, x + Inches(0.3), phase_y + Inches(1.8),
                phase_w - Inches(0.6), Inches(0.6),
                subtext, font_size=13, color=sub_color, line_spacing=1.25)

    # Arrow between phases
    if i < 2:
        arrow_x = x + phase_w + Inches(0.05)
        add_arrow_right(slide6, arrow_x, phase_y + Inches(1.0),
                        Inches(0.45), Inches(0.4), LGRAY)

# Metrics row
metric_y = Inches(4.8)
add_rect(slide6, Inches(0.8), metric_y - Inches(0.15), Inches(11.7), Inches(0.02), LGRAY)

add_textbox(slide6, Inches(0.8), metric_y, Inches(3), Inches(0.35),
            "Quality gates:", font_size=14, color=NAVY, bold=True,
            font_name=FONT_HEADING)

metrics = [
    ("100%", "factual accuracy\non financial details"),
    (">70%", "of drafts\nactually sent"),
    ("<20%", "of drafts\ndismissed"),
]

metric_x_start = Inches(0.8)
metric_w = Inches(3.6)
metric_inner_y = metric_y + Inches(0.5)

for i, (number, desc) in enumerate(metrics):
    x = metric_x_start + metric_w * i

    add_textbox(slide6, x, metric_inner_y, Inches(1.3), Inches(0.6),
                number, font_size=30, color=GOLD, bold=True,
                font_name=FONT_HEADING)

    add_textbox(slide6, x + Inches(1.3), metric_inner_y + Inches(0.05),
                Inches(2.2), Inches(0.6),
                desc, font_size=12, color=DGRAY, line_spacing=1.25)

# Silence callout
add_rounded_rect(slide6, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
                 RGBColor(0x08, 0x15, 0x1E))
add_textbox(slide6, Inches(1.2), Inches(6.15), Inches(11), Inches(0.5),
            "When AI isn\u2019t confident, it stays silent. "
            "Noise erodes trust faster than silence.",
            font_size=16, color=WHITE, bold=True)

# Measurement details
tf_meas = add_rich_textbox(slide6, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4))
add_para(tf_meas, "Every recommendation can be dismissed (tracked as a quality signal)  "
         "\u2022  Reasoning is transparent: pros see exactly which messages the AI used",
         font_size=11, color=MGRAY)


# ── SLIDE 7: "Where This Goes" ──────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
add_bg(slide7, NAVY)

add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
            "Where This Goes", font_size=44, color=WHITE, bold=True,
            font_name=FONT_HEADING)

# Bullet points with gold markers
bullets = [
    ("Memory grows with every conversation",
     "The AI gets better at understanding each customer and each business over time."),
    ("New skills plug into the same pipeline",
     "Estimate follow-ups, scheduling nudges, repeat customer recognition, warranty questions."),
    ("Architecture extends beyond invoicing",
     "The same pipeline could power AI assistance across other HCP workflows."),
    ("Earned autonomy as trust is proven",
     "Proactive outreach, cross-conversation insights, pattern detection across the customer base."),
]

bullet_y_start = Inches(1.7)
bullet_spacing = Inches(1.2)

for i, (title, desc) in enumerate(bullets):
    y = bullet_y_start + bullet_spacing * i

    # Gold marker circle
    marker = slide7.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.08), Inches(0.18), Inches(0.18)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = GOLD
    marker.line.fill.background()

    # Bullet title
    add_textbox(slide7, Inches(1.25), y - Inches(0.05), Inches(10.5), Inches(0.4),
                title, font_size=20, color=WHITE, bold=True)

    # Bullet description
    add_textbox(slide7, Inches(1.25), y + Inches(0.35), Inches(10.5), Inches(0.5),
                desc, font_size=15, color=MGRAY, line_spacing=1.3)

# Bottom callout
add_rect(slide7, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.02),
         RGBColor(0x1A, 0x2A, 0x3A))

add_textbox(slide7, Inches(0.8), Inches(6.45), Inches(11), Inches(0.6),
            "\u201cThis prototype is small on purpose. The point is to show "
            "the thinking and the foundation, not to ship everything at once.\u201d",
            font_size=16, color=GOLD, line_spacing=1.3)


# ── SLIDE 8: "Let Me Show You" ──────────────────────────────────
slide8 = prs.slides.add_slide(blank_layout)
add_bg(slide8, NAVY)

# Large centered title
add_textbox(slide8, Inches(0), Inches(1.0), SLIDE_W, Inches(1.0),
            "Let Me Show You", font_size=52, color=WHITE, bold=True,
            font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide8, Inches(0), Inches(2.0), SLIDE_W, Inches(0.5),
            "Three things to watch for:", font_size=20, color=GOLD,
            alignment=PP_ALIGN.CENTER)

# Three items centered
items = [
    ("1", "How the inbox surfaces what matters"),
    ("2", "How AI uses conversation history to draft relevant follow-ups"),
    ("3", "How a brand-new conversation gets evaluated in real time"),
]

item_y_start = Inches(3.2)
item_spacing = Inches(0.9)

for i, (num, text) in enumerate(items):
    y = item_y_start + item_spacing * i

    # Number in gold circle
    circ = slide8.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.2), y, Inches(0.5), Inches(0.5)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOLD
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = num
    circ.text_frame.paragraphs[0].font.size = Pt(18)
    circ.text_frame.paragraphs[0].font.color.rgb = NAVY
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].font.name = FONT_HEADING
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circ.text_frame.paragraphs[0].space_before = Pt(0)
    circ.text_frame.paragraphs[0].space_after = Pt(0)

    # Text
    add_textbox(slide8, Inches(4.9), y + Inches(0.02), Inches(5.5), Inches(0.5),
                text, font_size=18, color=WHITE, line_spacing=1.2)

# Bottom subtle line
add_rect(slide8, Inches(5.5), Inches(6.2), Inches(2.3), Inches(0.02), BLUE)


# ── Save ─────────────────────────────────────────────────────────
output_path = "/Users/briandesouza/Desktop/housecall-pro-ai-mvp/docs/HCP_AI_Smart_Inbox.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
