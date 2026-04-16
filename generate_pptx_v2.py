#!/usr/bin/env python3
"""Generate AI Smart Inbox v2 PowerPoint - 13 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
BG = RGBColor(0x08, 0x15, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xAA, 0xB8, 0xC2)
MGRAY = RGBColor(0x6B, 0x7B, 0x8D)
BLUE = RGBColor(0x00, 0x7A, 0xFF)
GREEN = RGBColor(0x34, 0xC7, 0x59)
ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
RED = RGBColor(0xFF, 0x3B, 0x30)
CARD = RGBColor(0x10, 0x20, 0x2E)
DARK = RGBColor(0x0A, 0x1E, 0x2E)

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
BL = prs.slide_layouts[6]

def bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG

def tb(s, l, t, w, h):
    f = s.shapes.add_textbox(l, t, w, h).text_frame
    f.word_wrap = True
    return f

def fp(f, txt, sz=14, c=WHITE, b=False, i=False, a=PP_ALIGN.LEFT):
    p = f.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.italic = i; r.font.name = 'Calibri'
    return p

def ap(f, txt, sz=14, c=WHITE, b=False, i=False, a=PP_ALIGN.LEFT, sb=0):
    p = f.add_paragraph(); p.alignment = a; p.space_before = Pt(sb)
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.italic = i; r.font.name = 'Calibri'
    return p

def rr(s, l, t, w, h, fill=CARD):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def rect(s, l, t, w, h, fill=CARD):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

# ============================================================
# SLIDE 1: I Started by Listening
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7)); fp(f, 'I Started by Listening', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4)); fp(f, 'Rauen Flooring  \u2022  Fort Myers, FL  \u2022  3 people', 13, MGRAY)

ph = rr(s, Inches(0.8), Inches(1.8), Inches(4.2), Inches(3.8))
f = ph.text_frame; f.word_wrap = True
p = f.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(70)
r = p.add_run(); r.text = '[ Family Photo ]'; r.font.size = Pt(16); r.font.color.rgb = MGRAY; r.font.name = 'Calibri'

quotes = [
    ('\u201cI\'ll finish a job covered in dust, check my phone, and there are six messages. The one about the invoice? Buried.\u201d', 'Sandro'),
    ('\u201cCustomers always ask \'why did you charge for leveling?\' We\'ve explained it a hundred times, but when we\'re on a job we can\'t stop and type it out.\u201d', 'Fabiana'),
    ('\u201cI thought Sandro answered her. He thought I did. She waited three days and hired someone else.\u201d', 'Fabiana'),
]
y = Inches(1.8)
for qt, who in quotes:
    f = tb(s, Inches(5.5), y, Inches(7.2), Inches(1.0))
    fp(f, qt, 13, LGRAY, i=True)
    ap(f, f'- {who}', 11, MGRAY, sb=4)
    y += Inches(1.3)

f = tb(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
fp(f, 'They were also clear about one thing: as customers themselves, they don\'t want to talk to a bot. They want to know a person is on the other end.', 12, MGRAY, i=True)

# ============================================================
# SLIDE 2: The Data Backs It Up
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'The Data Backs It Up', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5)); fp(f, 'What I heard at the kitchen table is not just my family\'s problem.', 16, LGRAY)

stats = [('56%', 'of small businesses are owed\nmoney from unpaid invoices'),
         ('$17.5K', 'average amount owed\nper business'),
         ('94% \u2192 50%', 'collection probability drops\nfrom 30 days to 6 months'),
         ('+50%', 'payment likelihood increase\nwith just one reminder')]
bw = Inches(2.6); gap = Inches(0.3)
for i, (st, lb) in enumerate(stats):
    x = Inches(0.8) + i * (bw + gap)
    rr(s, x, Inches(2.5), bw, Inches(2.5))
    f = tb(s, x + Inches(0.2), Inches(2.8), bw - Inches(0.4), Inches(0.8)); fp(f, st, 36, BLUE, b=True, a=PP_ALIGN.CENTER)
    f = tb(s, x + Inches(0.2), Inches(3.7), bw - Inches(0.4), Inches(1.0)); fp(f, lb, 12, LGRAY, a=PP_ALIGN.CENTER)

f = tb(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5)); fp(f, 'Speed and clarity are the two biggest levers for getting paid.', 16, WHITE, b=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Where the Funnel Leaks
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'Where the Funnel Leaks', 36, b=True)

stages = ['Invoice\nSent', 'Customer\nReceives', 'Questions &\nConcerns', 'Pro\nResponds', 'Payment\nMade']
sw = Inches(1.8); sg = Inches(0.35)
for i, st in enumerate(stages):
    x = Inches(1.2) + i * (sw + sg)
    c = RED if i == 2 else CARD
    bx = rr(s, x, Inches(1.5), sw, Inches(1.0), c)
    f = bx.text_frame; f.word_wrap = True
    p = f.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = st; r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = 'Calibri'

f = tb(s, Inches(1.2) + 2*(sw+sg), Inches(2.6), sw, Inches(0.4)); fp(f, '\u2191 THE GAP', 11, RED, b=True, a=PP_ALIGN.CENTER)

f = tb(s, Inches(0.8), Inches(3.2), Inches(5.2), Inches(0.4)); fp(f, 'What HCP does well today', 14, GREEN, b=True)
f = tb(s, Inches(0.8), Inches(3.7), Inches(5.2), Inches(1.5))
fp(f, '\u2022  Invoice creation & multiple payment methods', 12, LGRAY)
ap(f, '\u2022  Automated email reminders', 12, LGRAY, sb=4)
ap(f, '\u2022  Marketing AI for invoice message templates', 12, LGRAY, sb=4)

f = tb(s, Inches(6.8), Inches(3.2), Inches(5.8), Inches(0.4)); fp(f, 'What competitors are adding', 14, ORANGE, b=True)
f = tb(s, Inches(6.8), Inches(3.7), Inches(5.8), Inches(1.5))
fp(f, '\u2022  ServiceTitan: AI invoice summaries + one-click collection emails', 12, LGRAY)
ap(f, '\u2022  Jobber: AI rewrite for invoice messages, reminders, two-way texting', 12, LGRAY, sb=4)
ap(f, '\u2022  Workiz: AI messaging across channels, tailored to customer history', 12, LGRAY, sb=4)

rr(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2))
f = tb(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.3)); fp(f, 'THE OPPORTUNITY', 12, BLUE, b=True)
f = tb(s, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.6)); fp(f, 'Competitors are adding AI to individual touchpoints. None have a unified AI layer that watches the full conversation, understands context from memory, and helps the pro respond at the moment it matters most.', 13, LGRAY)

# ============================================================
# SLIDE 4: AI Smart Inbox
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'AI Smart Inbox', 40, b=True)
f = tb(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5)); fp(f, 'An AI co-worker that lives inside HouseCall Pro\'s messaging inbox.', 18, LGRAY)

caps = [('PRIORITIZE', 'Surface what needs your attention right now. No hunting through threads.'),
        ('BE QUICK', 'AI-drafted responses ready the moment you check your inbox. Turn days into minutes.'),
        ('DRIVE CLARITY', 'Use context that already exists (conversation history, invoice details, your own FAQ answers) to draft accurate, helpful replies.')]
cw = Inches(3.6); cg = Inches(0.3)
for i, (t, d) in enumerate(caps):
    x = Inches(0.8) + i * (cw + cg)
    rr(s, x, Inches(2.5), cw, Inches(2.8))
    f = tb(s, x + Inches(0.25), Inches(2.7), cw - Inches(0.5), Inches(0.5)); fp(f, t, 16, BLUE, b=True)
    f = tb(s, x + Inches(0.25), Inches(3.3), cw - Inches(0.5), Inches(1.8)); fp(f, d, 13, LGRAY)

f = tb(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8))
fp(f, 'Design principle: AI works behind the scenes. The pro reviews, edits, and sends every message. No customer ever talks to a bot.', 13, MGRAY, i=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: What It Looks Like
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'What It Looks Like', 36, b=True)

panels = [('The Inbox', 'Pending Actions cards at the top surface the 2-3 things that need attention. Priority indicators help the pro triage at a glance. All other conversations stay accessible below.', 'AI surfaces what matters.\nEverything else is one tap away.'),
          ('Conversation + Draft', 'When the pro taps into an action, the AI draft appears at the bottom of the chat thread. An editable textarea with a Send button. The draft sits alongside the real conversation so the pro can verify it.', 'Draft appears in context.\nReview, edit, send.'),
          ('AI Reasoning', 'Every draft has a "Why this?" button. It opens the full evidence chain: source messages with timestamps, business knowledge entries, the reasoning chain, and a confidence score.', 'Full transparency.\nThe pro sees the AI\'s homework.')]
pw = Inches(3.6); pg = Inches(0.3)
for i, (t, d, an) in enumerate(panels):
    x = Inches(0.8) + i * (pw + pg)
    ph = rr(s, x, Inches(1.5), pw, Inches(2.8))
    f = ph.text_frame; f.word_wrap = True
    p = f.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(50)
    r = p.add_run(); r.text = f'[ {t} Screenshot ]'; r.font.size = Pt(14); r.font.color.rgb = MGRAY; r.font.name = 'Calibri'
    f = tb(s, x, Inches(4.5), pw, Inches(0.4)); fp(f, t, 14, WHITE, b=True)
    f = tb(s, x, Inches(4.9), pw, Inches(1.0)); fp(f, d, 10, LGRAY)
    f = tb(s, x, Inches(6.1), pw, Inches(0.6)); fp(f, an, 10, BLUE, i=True)

# ============================================================
# SLIDE 6: The Four Skills
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'The Four Skills', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4)); fp(f, 'Each skill has its own prompt, its own tone, and its own trigger criteria.', 14, LGRAY)

skills = [
    ('Invoice Question', RED, 'HIGH', 'Customer questions a specific charge', '"What\'s the float charge? $450 seems like a lot."', 'Explains the charge using invoice data, FAQ, and prior conversation context.'),
    ('Invoice Summary', GREEN, 'LOW', 'New invoice sent without a friendly breakdown', '(Proactive: no customer message needed)', 'Breaks down each line item, explains value, mentions payment options.'),
    ('Payment Follow-up', ORANGE, 'HIGH', 'Customer promised to pay by a date that has passed', '"I\'ll pay by Friday" (it\'s now Monday)', 'Friendly nudge referencing the customer\'s own words. Not a demand.'),
    ('Payment Plan', BLUE, 'HIGH', 'Customer expresses difficulty paying the full amount', '"$5,800 is a lot. Can I break this up?"', 'Understanding tone. Suggests 2-3 installments with specific amounts.'),
]
rh = Inches(1.05); sy = Inches(1.9)
for i, (nm, clr, pri, trg, ex, appr) in enumerate(skills):
    y = sy + i * (rh + Inches(0.15))
    rr(s, Inches(0.8), y, Inches(11.7), rh)
    f = tb(s, Inches(1.0), y + Inches(0.1), Inches(2.2), Inches(0.35)); fp(f, nm, 14, clr, b=True)
    f = tb(s, Inches(1.0), y + Inches(0.45), Inches(2.2), Inches(0.25)); fp(f, pri, 9, clr, b=True)
    f = tb(s, Inches(3.3), y + Inches(0.1), Inches(3.5), Inches(0.35)); fp(f, trg, 11, WHITE)
    f = tb(s, Inches(3.3), y + Inches(0.5), Inches(3.5), Inches(0.45)); fp(f, ex, 10, MGRAY, i=True)
    f = tb(s, Inches(7.0), y + Inches(0.15), Inches(5.3), Inches(0.8)); fp(f, appr, 11, LGRAY)

f = tb(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5))
fp(f, 'When AI isn\'t confident, it stays silent. The demo includes conversations with no AI action, because none was needed.', 12, MGRAY, i=True)

# ============================================================
# SLIDE 7: Under the Hood
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'Under the Hood', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4)); fp(f, 'Two-stage pipeline with a persistent memory layer', 16, LGRAY)

# Stage 1
rr(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2))
f = tb(s, Inches(1.0), Inches(2.1), Inches(5.1), Inches(0.4)); fp(f, '1  EVALUATE', 16, BLUE, b=True)
f = tb(s, Inches(1.0), Inches(2.6), Inches(5.1), Inches(1.4))
fp(f, 'Every incoming message runs through an AI agent that decides:', 12, LGRAY)
ap(f, '\u2022  Does this need action?', 12, LGRAY, sb=6)
ap(f, '\u2022  What kind? How urgent?', 12, LGRAY, sb=4)
ap(f, '\u2022  Update customer memory', 12, LGRAY, sb=4)

# Stage 2
rr(s, Inches(6.6), Inches(2.0), Inches(5.9), Inches(2.2))
f = tb(s, Inches(6.8), Inches(2.1), Inches(5.5), Inches(0.4)); fp(f, '2  GENERATE', 16, BLUE, b=True)
f = tb(s, Inches(6.8), Inches(2.6), Inches(5.5), Inches(1.4))
fp(f, 'If action needed, a second pass drafts a response using:', 12, LGRAY)
ap(f, '\u2022  Invoice details', 12, LGRAY, sb=6)
ap(f, '\u2022  Conversation history', 12, LGRAY, sb=4)
ap(f, '\u2022  Two layers of memory', 12, LGRAY, sb=4)

# Memory layer
f = tb(s, Inches(0.8), Inches(4.5), Inches(2.0), Inches(0.4)); fp(f, 'MEMORY LAYER', 14, GREEN, b=True)

rr(s, Inches(0.8), Inches(5.0), Inches(5.8), Inches(0.9))
f = tb(s, Inches(1.0), Inches(5.05), Inches(5.4), Inches(0.8))
fp(f, 'Business Knowledge', 13, WHITE, b=True)
ap(f, 'Your services, pricing, FAQ answers. Things you explain every week.', 11, LGRAY, sb=4)

rr(s, Inches(6.85), Inches(5.0), Inches(5.65), Inches(0.9))
f = tb(s, Inches(7.05), Inches(5.05), Inches(5.25), Inches(0.8))
fp(f, 'Customer History', 13, WHITE, b=True)
ap(f, 'Preferences, past jobs, communication style. Grows with every interaction.', 11, LGRAY, sb=4)

# Bottom line - pure mechanics, no forward-looking
f = tb(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
fp(f, 'Memory is the differentiator. The AI doesn\'t treat each message in isolation. It builds understanding of each customer over time, just like the pro does.', 13, WHITE, b=True)

# ============================================================
# SLIDE 8: Trust is Earned
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'Trust is Earned', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4)); fp(f, 'Progression is earned by metrics, not by a calendar.', 16, LGRAY)

phases = [('PHASE 1  \u2022  This MVP', 'Every draft reviewed\nby the pro', 'Building confidence', BLUE),
          ('PHASE 2  \u2022  Next', 'Auto-send for proven,\nlow-risk actions', 'Track record earns automation', GREEN),
          ('PHASE 3  \u2022  Future', 'Autonomous for\nroutine cases', 'AI operates where trust\nis established', ORANGE)]
phw = Inches(3.6); phg = Inches(0.3)
for i, (t, beh, out, clr) in enumerate(phases):
    x = Inches(0.8) + i * (phw + phg)
    rr(s, x, Inches(1.9), phw, Inches(1.6))
    f = tb(s, x + Inches(0.2), Inches(2.0), phw - Inches(0.4), Inches(0.35)); fp(f, t, 12, clr, b=True)
    f = tb(s, x + Inches(0.2), Inches(2.4), phw - Inches(0.4), Inches(0.6)); fp(f, beh, 13, WHITE)
    f = tb(s, x + Inches(0.2), Inches(3.05), phw - Inches(0.4), Inches(0.4)); fp(f, out, 11, MGRAY, i=True)

# Quality gates
f = tb(s, Inches(0.8), Inches(3.7), Inches(3.0), Inches(0.4)); fp(f, 'Quality gates:', 13, WHITE, b=True)
gates = [('100%', 'factual accuracy\non financial details'), ('>70%', 'of drafts\nactually sent'), ('<20%', 'of drafts\ndismissed')]
for i, (st, lb) in enumerate(gates):
    x = Inches(0.8) + i * Inches(2.5)
    f = tb(s, x, Inches(4.1), Inches(2.2), Inches(0.5)); fp(f, st, 28, BLUE, b=True, a=PP_ALIGN.CENTER)
    f = tb(s, x, Inches(4.7), Inches(2.2), Inches(0.5)); fp(f, lb, 10, LGRAY, a=PP_ALIGN.CENTER)

# AI vs Human
f = tb(s, Inches(8.0), Inches(3.7), Inches(4.5), Inches(0.4)); fp(f, 'AI vs. Human', 13, WHITE, b=True)
aih = [('AI scans, drafts, prioritizes', 'Human decides, edits, sends'),
       ('AI suggests', 'Human approves'),
       ('AI stays silent when unsure', 'Human owns the relationship')]
f = tb(s, Inches(8.0), Inches(4.1), Inches(4.5), Inches(1.5))
for i, (ai, hu) in enumerate(aih):
    if i == 0:
        fp(f, f'{ai}  |  {hu}', 10, LGRAY)
    else:
        ap(f, f'{ai}  |  {hu}', 10, LGRAY, sb=6)

f = tb(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4)); fp(f, 'When AI isn\'t confident, it stays silent. Noise erodes trust faster than silence.', 13, WHITE, b=True)
f = tb(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5))
fp(f, 'Every recommendation can be dismissed (tracked as a quality signal)  \u2022  Reasoning is transparent: pros see exactly which messages the AI used', 11, MGRAY)

# ============================================================
# SLIDE 9: How We Measure
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'How We Measure', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4)); fp(f, 'Metrics that prove it works', 16, LGRAY)

# Left column
rr(s, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.4))
f = tb(s, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.4)); fp(f, 'Outcome Metrics', 15, GREEN, b=True)
f = tb(s, Inches(1.0), Inches(2.4), Inches(5.1), Inches(0.3)); fp(f, 'Does it move the business needle?', 11, MGRAY, i=True)
oms = [('Days-to-payment', 'Decrease by 20%+'), ('Pro response rate (within 24h)', '40% \u2192 80%+'),
       ('30-day collection rate', 'Increase by 15%+'), ('Revenue recovered', 'Track absolute $'),
       ('Pro time saved', 'Decrease by 50%+')]
y = Inches(2.9)
for m, t in oms:
    f = tb(s, Inches(1.0), y, Inches(3.2), Inches(0.3)); fp(f, m, 12, WHITE)
    f = tb(s, Inches(4.2), y, Inches(1.8), Inches(0.3)); fp(f, t, 12, BLUE, b=True)
    y += Inches(0.38)

# Right column
rr(s, Inches(6.7), Inches(1.9), Inches(5.8), Inches(4.4))
f = tb(s, Inches(6.9), Inches(2.0), Inches(5.4), Inches(0.4)); fp(f, 'Quality Metrics', 15, BLUE, b=True)
f = tb(s, Inches(6.9), Inches(2.4), Inches(5.4), Inches(0.3)); fp(f, 'Is the AI doing a good job?', 11, MGRAY, i=True)
qms = [('Draft send rate', '> 70%'), ('Draft acceptance (no edits)', '> 50%'),
       ('Dismissal rate', '< 20%'), ('Factual accuracy', '100%'), ('Pro NPS on feature', '> 30')]
y = Inches(2.9)
for m, t in qms:
    f = tb(s, Inches(6.9), y, Inches(3.5), Inches(0.3)); fp(f, m, 12, WHITE)
    f = tb(s, Inches(10.5), y, Inches(1.8), Inches(0.3)); fp(f, t, 12, BLUE, b=True)
    y += Inches(0.38)

f = tb(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
fp(f, '100% factual accuracy is a hard requirement. A single wrong dollar amount in a customer message could damage the pro\'s credibility.', 12, RED, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: Quality Gates
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'Quality Gates', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4)); fp(f, 'We don\'t graduate on a calendar. We graduate on metrics.', 16, LGRAY)

gd = [('Internal \u2192 Beta', 'No critical bugs; team consensus on draft quality'),
      ('Beta \u2192 A/B Test', 'Send rate > 60%; dismiss < 25%; zero factual errors; NPS > 20'),
      ('A/B \u2192 General Availability', 'Significant improvement in days-to-payment; send rate > 70%; dismiss < 20%; NPS > 30'),
      ('Phase 1 \u2192 Phase 2', '90+ days sustained quality at scale; explicit pro demand; zero customer complaints')]
y = Inches(2.0)
for g, c in gd:
    rr(s, Inches(0.8), y, Inches(11.7), Inches(0.9))
    f = tb(s, Inches(1.0), y + Inches(0.15), Inches(3.0), Inches(0.6)); fp(f, g, 14, BLUE, b=True)
    f = tb(s, Inches(4.2), y + Inches(0.15), Inches(8.1), Inches(0.6)); fp(f, c, 12, LGRAY)
    y += Inches(1.1)

rr(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.9), RGBColor(0x2A, 0x15, 0x15))
f = tb(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.7))
fp(f, 'If any guardrail is breached (dismiss rate > 20%, any factual error, CSAT drops > 5%), we pause and investigate before continuing.', 13, RED, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: Who It's For
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'Who It\'s For', 36, b=True)

personas = [
    ('\u201cMike\u201d', 'Solo Plumber', '$180K/year  \u2022  4-6 jobs/day',
     '$12K in unpaid invoices. Sends one follow-up, then drops it. Hates feeling like a bill collector. Payment messages get lost in scheduling chatter.', BLUE),
    ('\u201cMaria\u201d', '12-Employee Cleaning Company', '$450K/year  \u2022  80-120 jobs/month',
     '15-25% of invoices overdue at any time. Nobody "owns" follow-up. Customer questions sit for days because the office helper can\'t answer billing questions.', GREEN),
    ('\u201cJesse\u201d', '4-Person HVAC Company', '$600K/year  \u2022  $2K-$8K jobs',
     '$35K outstanding. Only Jesse can explain technical pricing. He\'s on a roof for 2-3 days before he can respond to "Why is this $800 more than the estimate?"', ORANGE),
]
cw = Inches(3.6); cg = Inches(0.3)
for i, (nm, role, st, pain, clr) in enumerate(personas):
    x = Inches(0.8) + i * (cw + cg)
    rr(s, x, Inches(1.5), cw, Inches(4.0))
    f = tb(s, x + Inches(0.25), Inches(1.7), cw - Inches(0.5), Inches(0.5)); fp(f, nm, 22, clr, b=True)
    f = tb(s, x + Inches(0.25), Inches(2.2), cw - Inches(0.5), Inches(0.4)); fp(f, role, 13, WHITE, b=True)
    f = tb(s, x + Inches(0.25), Inches(2.6), cw - Inches(0.5), Inches(0.3)); fp(f, st, 10, MGRAY)
    rect(s, x + Inches(0.25), Inches(3.0), cw - Inches(0.5), Inches(0.02), clr)
    f = tb(s, x + Inches(0.25), Inches(3.2), cw - Inches(0.5), Inches(2.0)); fp(f, pain, 12, LGRAY)

f = tb(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.5))
fp(f, 'Different businesses, same problem: payment-critical messages get buried, and every day of silence costs money.', 14, WHITE, b=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: The Bigger Picture (merged vision + future)
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7)); fp(f, 'The Bigger Picture', 36, b=True)
f = tb(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4))
fp(f, 'This MVP is the starting point. The architecture is built for more.', 16, LGRAY)

futures = [
    ('Memory compounds over time', 'The AI gets better at understanding each customer and each business with every conversation. Six months of usage produces a richer, more accurate assistant.'),
    ('New skills plug into the same pipeline', 'Estimate follow-ups, scheduling nudges, repeat customer recognition, warranty questions. Each new skill is a prompt, not a rebuild.'),
    ('Multiple entry points', 'Messages today. Emails and notifications tomorrow. Every new entry point adds context, and richer context enables more capable AI.'),
    ('Earned autonomy as trust is proven', 'Proactive outreach, cross-conversation insights, pattern detection across the customer base.'),
]
y = Inches(1.9)
for i, (t, d) in enumerate(futures):
    rr(s, Inches(0.8), y, Inches(11.7), Inches(0.95))
    f = tb(s, Inches(1.0), y + Inches(0.12), Inches(0.5), Inches(0.5)); fp(f, str(i+1), 22, BLUE, b=True)
    f = tb(s, Inches(1.5), y + Inches(0.08), Inches(10.7), Inches(0.4)); fp(f, t, 15, WHITE, b=True)
    f = tb(s, Inches(1.5), y + Inches(0.48), Inches(10.7), Inches(0.4)); fp(f, d, 12, LGRAY)
    y += Inches(1.1)

# Key insight box
rr(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.7), DARK)
f = tb(s, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.6))
fp(f, 'The richer the context, the more the AI can do. HouseCall Pro already has the data. This MVP shows what happens when you give an AI layer access to it.', 13, WHITE, b=True, a=PP_ALIGN.CENTER)

f = tb(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4))
fp(f, '"This prototype is small on purpose. The point is to show the thinking and the foundation, not to ship everything at once."', 12, MGRAY, i=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: Let Me Show You
# ============================================================
s = prs.slides.add_slide(BL); bg(s)
f = tb(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0)); fp(f, 'Let Me Show You', 48, b=True, a=PP_ALIGN.CENTER)
f = tb(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.5)); fp(f, 'Three things to watch for:', 16, LGRAY, a=PP_ALIGN.CENTER)

watch = [('1', 'How the inbox surfaces what matters'),
         ('2', 'How AI uses conversation history and business knowledge to draft relevant responses'),
         ('3', 'How a new message gets evaluated live, and when AI chooses to stay silent')]
y = Inches(3.6)
for n, t in watch:
    c = rr(s, Inches(3.5), y, Inches(0.5), Inches(0.5), BLUE)
    f = c.text_frame
    p = f.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n; r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = 'Calibri'
    f = tb(s, Inches(4.3), y, Inches(7.0), Inches(0.5)); fp(f, t, 15, WHITE)
    y += Inches(0.8)

# ============================================================
prs.save('docs/HCP_AI_Smart_Inbox_v2.pptx')
print(f'Saved docs/HCP_AI_Smart_Inbox_v2.pptx with {len(prs.slides)} slides')
